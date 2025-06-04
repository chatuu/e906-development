from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
# from pptx.dml.color import RGBColor # Not used, but good for color customization
from datetime import datetime
import os
import platform # To check OS

# --- Configuration ---
AUTHOR_NAME = "Chatura Kuruppu"
UNIVERSITY_NAME = "New Mexico State University"
MAIN_TITLE_TEXT = "Weekly Updates"
OVERVIEW_TITLE_TEXT = "Overview"

# Bullet points for the Overview slide
OVERVIEW_BULLETS = [
    "This is the first bullet point for the overview.",
    "This is the second bullet point, providing more details.",
    "And finally, the third bullet point summarizing key aspects."
]

# Font sizes
MAIN_TITLE_FONT_SIZE = Pt(48)
SUB_TITLE_FONT_SIZE = Pt(24)
SLIDE_TITLE_FONT_SIZE = Pt(36)  # For "Overview" title
BODY_FONT_SIZE = Pt(18)
FOOTER_FONT_SIZE = Pt(10)

# Font names
FONT_TITLE = 'Calibri Light'
FONT_BODY = 'Calibri'
FONT_FOOTER = 'Arial'

# --- Helper Functions ---
def get_formatted_date():
    """Returns the current date formatted as YYYY-MM-DD."""
    return datetime.now().strftime("%Y-%m-%d")

def add_footer_to_slide(slide, date_str, prs_width, prs_height):
    """
    Adds a three-part footer to the given slide.
    Assumes slide dimensions are passed via prs_width and prs_height.
    """
    footer_height = Inches(0.3)
    bottom_margin = Inches(0.1)
    side_margin = Inches(0.25)
    
    footer_top = prs_height - footer_height - bottom_margin
    footer_box_width = Inches(3.0) 

    # Left Footer
    left_ft_left = side_margin
    left_ft_box = slide.shapes.add_textbox(left_ft_left, footer_top, 
                                           footer_box_width, footer_height)
    left_tf = left_ft_box.text_frame
    left_tf.clear() 
    p_left = left_tf.add_paragraph()
    p_left.text = AUTHOR_NAME
    p_left.font.size = FOOTER_FONT_SIZE
    p_left.font.name = FONT_FOOTER
    p_left.alignment = PP_ALIGN.LEFT
    left_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Center Footer
    center_ft_left = (prs_width / 2) - (footer_box_width / 2)
    center_ft_box = slide.shapes.add_textbox(center_ft_left, footer_top,
                                             footer_box_width, footer_height)
    center_tf = center_ft_box.text_frame
    center_tf.clear()
    p_center = center_tf.add_paragraph()
    p_center.text = UNIVERSITY_NAME
    p_center.font.size = FOOTER_FONT_SIZE
    p_center.font.name = FONT_FOOTER
    p_center.alignment = PP_ALIGN.CENTER
    center_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Right Footer
    right_ft_left = prs_width - side_margin - footer_box_width
    right_ft_box = slide.shapes.add_textbox(right_ft_left, footer_top,
                                            footer_box_width, footer_height)
    right_tf = right_ft_box.text_frame
    right_tf.clear()
    p_right = right_tf.add_paragraph()
    p_right.text = date_str
    p_right.font.size = FOOTER_FONT_SIZE
    p_right.font.name = FONT_FOOTER
    p_right.alignment = PP_ALIGN.RIGHT
    right_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def save_as_pdf_windows(pptx_path, pdf_path):
    """
    Saves a PowerPoint presentation as PDF using COM automation on Windows.
    Requires Microsoft PowerPoint to be installed and the 'comtypes' library.
    """
    if platform.system() == "Windows":
        try:
            import comtypes.client
            
            powerpoint = None  # Initialize to ensure it's defined in finally block
            presentation = None

            # Get absolute paths
            abs_pptx_path = os.path.abspath(pptx_path)
            abs_pdf_path = os.path.abspath(pdf_path)

            try:
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                # Make PowerPoint invisible during conversion
                # powerpoint.Visible = 1 # msoTrue, use for debugging if needed
                
                # Open the presentation
                presentation = powerpoint.Presentations.Open(abs_pptx_path)
                
                # Save as PDF (formatType 32 for PDF)
                presentation.SaveAs(abs_pdf_path, 32) 
                print(f"Presentation successfully saved as PDF: {abs_pdf_path}")
            except Exception as e:
                print(f"Error during PDF conversion: {e}")
                print("Please ensure Microsoft PowerPoint is installed and 'comtypes' library is working.")
            finally:
                if presentation:
                    presentation.Close()
                if powerpoint:
                    powerpoint.Quit()
                # Ensure comtypes releases the COM objects
                comtypes.CoUninitialize()

        except ImportError:
            print("Skipping PDF conversion: 'comtypes' library not found. "
                  "Install it with 'pip install comtypes' (Windows only).")
        except Exception as e:
            print(f"An unexpected error occurred while trying to import or use comtypes: {e}")
    else:
        print(f"Skipping PDF conversion: This feature is currently supported on Windows only.")

# --- Main Script ---
def create_presentation():
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    current_date_str = get_formatted_date()

    # === Slide 1: Title Slide ===
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)

    title_shape = slide1.shapes.title
    tf_title = title_shape.text_frame
    if not tf_title.paragraphs:
        p_title = tf_title.add_paragraph()
    else:
        p_title = tf_title.paragraphs[0]
        p_title.clear()
    
    run_title = p_title.add_run()
    run_title.text = MAIN_TITLE_TEXT
    run_title.font.size = MAIN_TITLE_FONT_SIZE
    run_title.font.name = FONT_TITLE
    p_title.alignment = PP_ALIGN.CENTER

    try:
        subtitle_shape = slide1.placeholders[1] 
        tf_subtitle = subtitle_shape.text_frame
        if not tf_subtitle.paragraphs:
            p_subtitle = tf_subtitle.add_paragraph()
        else:
            p_subtitle = tf_subtitle.paragraphs[0]
            p_subtitle.clear()
        
        run_subtitle = p_subtitle.add_run()
        run_subtitle.text = current_date_str
        run_subtitle.font.size = SUB_TITLE_FONT_SIZE
        run_subtitle.font.name = FONT_TITLE
        p_subtitle.alignment = PP_ALIGN.CENTER
    except KeyError:
        print("Warning: Subtitle placeholder (index 1) not found on title slide layout.")
        sub_title_box = slide1.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1))
        tf_subtitle_manual = sub_title_box.text_frame
        p_subtitle_manual = tf_subtitle_manual.add_paragraph()
        p_subtitle_manual.text = current_date_str
        p_subtitle_manual.font.size = SUB_TITLE_FONT_SIZE
        p_subtitle_manual.font.name = FONT_TITLE
        p_subtitle_manual.alignment = PP_ALIGN.CENTER

    add_footer_to_slide(slide1, current_date_str, slide_width, slide_height)

    # === Slide 2: First Slide (Overview) ===
    blank_slide_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(blank_slide_layout)

    title_width_overview = Inches(4.5)
    title_height_overview = Inches(0.75)
    title_top_overview = Inches(0.5)
    title_left_overview = slide_width - title_width_overview - Inches(0.5) 

    overview_title_box = slide2.shapes.add_textbox(title_left_overview, title_top_overview, 
                                                   title_width_overview, title_height_overview)
    tf_overview_title = overview_title_box.text_frame
    tf_overview_title.clear()
    p_overview_title = tf_overview_title.add_paragraph()
    p_overview_title.text = OVERVIEW_TITLE_TEXT
    p_overview_title.font.size = SLIDE_TITLE_FONT_SIZE
    p_overview_title.font.name = FONT_BODY
    p_overview_title.font.bold = True
    p_overview_title.alignment = PP_ALIGN.RIGHT
    tf_overview_title.vertical_anchor = MSO_ANCHOR.TOP 

    body_top_overview = title_top_overview + title_height_overview + Inches(0.3) 
    body_left_overview = Inches(0.75)
    body_width_overview = slide_width - Inches(1.5) 
    body_height_overview = Inches(4) 

    body_box = slide2.shapes.add_textbox(body_left_overview, body_top_overview, 
                                         body_width_overview, body_height_overview)
    tf_body = body_box.text_frame
    tf_body.clear()
    tf_body.word_wrap = True

    for item_text in OVERVIEW_BULLETS:
        p_item = tf_body.add_paragraph()
        p_item.text = item_text
        p_item.font.size = BODY_FONT_SIZE
        p_item.font.name = FONT_BODY
        p_item.level = 0 
        p_item.space_after = Pt(6)
    
    add_footer_to_slide(slide2, current_date_str, slide_width, slide_height)

    # --- Save Presentation ---
    base_file_name = f"Weekly_Updates_{current_date_str}"
    pptx_file_name = f"{base_file_name}.pptx"
    pdf_file_name = f"{base_file_name}.pdf"

    try:
        prs.save(pptx_file_name)
        print(f"Presentation '{pptx_file_name}' created successfully.")
        
        # Attempt to save as PDF (Windows with PowerPoint only)
        save_as_pdf_windows(pptx_file_name, pdf_file_name)
        
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == '__main__':
    create_presentation()